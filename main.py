# type: ignore

import os
import PyPDF2
import fitz  # PyMuPDF
import pdfplumber
import pandas as pd
from pptx import Presentation
from docx import Document as DocxDocument

from dotenv import load_dotenv
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_experimental.text_splitter import SemanticChunker
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_core.prompts import PromptTemplate
from langchain.chains import ConversationalRetrievalChain
from langchain_groq import ChatGroq
from langchain.memory import ConversationBufferMemory
from langchain_core.documents import Document

from interface.ui import build_interface

# 🔑 Load API Key
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

# ------------------------
# Configurable Model Choices
# ------------------------
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
GROQ_MODEL = "llama-3.3-70b-versatile"

vectorstore = None
qa_chain = None

# ------------------------
# Translation Helpers
# ------------------------
def translate_to_urdu(text: str) -> str:
    # simple placeholder - replace with real translation model or API if needed
    return f"[اردو ترجمہ]: {text}"

def translate_from_urdu(question: str) -> str:
    # placeholder for Urdu → English if needed
    return f"(Translated from Urdu) {question}"

# ------------------------
# PDF Extraction Utilities
# ------------------------
def extract_pdf_content(file_path):      # this function read the files n three different method way and extract the text
    content = ""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)             # pypdf2
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text and text.strip():
                    content += f"Page {i+1}:\n{text}\n\n"
        if content.strip():
            return content
    except:
        pass

    try:
        doc = fitz.open(file_path)                 # pymudpdf
        for i in range(len(doc)):
            text = doc[i].get_text()
            if text and text.strip():
                content += f"Page {i+1}:\n{text}\n\n"
        doc.close()
        if content.strip():
            return content
    except:
        pass

    try:
        with pdfplumber.open(file_path) as pdf:        # pdfplumber
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text and text.strip():
                    content += f"Page {i+1}:\n{text}\n\n"
        if content.strip():
            return content
    except:
        pass

    return None

# ------------------------
# Load Documents
# ------------------------
def load_documents(files):
    docs = []
    for file in files:
        file_path = file.name if hasattr(file, "name") else file

        if file_path.endswith(".pdf"):
            content = extract_pdf_content(file_path)
            if content:
                docs.append(Document(page_content=content, metadata={"source": file_path}))
            else:
                docs.extend(PyPDFLoader(file_path).load())

        elif file_path.endswith(".docx"):
            docs.extend(Docx2txtLoader(file_path).load())

        elif file_path.endswith(".txt"):
            docs.extend(TextLoader(file_path).load())

        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            if text.strip():
                docs.append(Document(page_content=text, metadata={"source": file_path}))

        elif file_path.endswith(".csv"):
            df = pd.read_csv(file_path)
            docs.append(Document(page_content=df.to_string(), metadata={"source": file_path}))

        elif file_path.endswith(".xlsx"):
            df = pd.read_excel(file_path)
            docs.append(Document(page_content=df.to_string(), metadata={"source": file_path}))

    return docs

# ------------------------
# Process Documents
# ------------------------
def process_documents(docs):
    embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

    try:
        splitter = SemanticChunker(embeddings)
        splits = splitter.split_documents(docs)
    except:
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200, add_start_index=True)
        splits = splitter.split_documents(docs)

    # Add page number if available
    for i, doc in enumerate(splits):
        if "page" not in doc.metadata:
            doc.metadata["page"] = i + 1   # agar nahi hai toh assign

    return FAISS.from_documents(splits, embeddings)

# ------------------------
# Create Chain
# ------------------------
def create_chain(vectorstore):    # This function builds the conversation system (chain) where:  4 thing 1 llm 2 prompt template 3 vectorstore 4 memory
    llm = ChatGroq(
        api_key=GROQ_API_KEY,
        model=GROQ_MODEL,
        temperature=0.5,
        max_tokens=1024
    )

      # 🔹 Smart, general-purpose assistant prompt                      # your rules and instructions.
    prompt_template = """                                                                      
You are a highly intelligent and professional AI assistant. 
Your role is to help users explore and understand their uploaded documents in the most effective way.

## Instructions:
1. *Answer strictly based on the user's request.*
   - If the user asks for definition only → provide only the definition.
   - If the user asks for examples only → provide only examples.
   - If the user asks for key points, characteristics, advantages, disadvantages, steps, or summary → provide exactly that.
2. *Be adaptive like an expert tutor*:
   - If the question is vague → provide a clear, structured explanation.
   - If the question is precise → provide a direct, concise answer.
3. *Never hallucinate*:
   - If the answer is not found in the documents, reply:  
     "This information is not available in the uploaded documents."
4. *Formatting*:
   - Use *markdown* for readability.
   - Prefer *headings, bullet points, and tables* where helpful.
   - Keep tone *professional, user-friendly, and easy to follow*.

## Context:                       
{context}            #   Context:  (retrieved chunks here)

## Question:
{question}           # Question:   (users query here)

## Answer:
"""
 
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])  

    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)    # remembers past conversation.

    chain = ConversationalRetrievalChain.from_llm(                        # It’s a LangChain class.that  Connects LLM + Retriever (FAISS) + Prompt + Memory.
        llm,
        retriever=vectorstore.as_retriever(search_kwargs={"k": 8}),   # retriver automatically so uery query embiding and plus find simlierty of vectorstore
        memory=memory,
        combine_docs_chain_kwargs={"prompt": prompt}    # is prompt ma query embedding and chunks chain ke ander khud huta hai
    )
    return chain


# ------------------------
# Chatbot Response
# ------------------------
def chatbot_response(files, question, language):
    global vectorstore, qa_chain

    # Load docs if first time
    if files and vectorstore is None:
        docs = load_documents(files)
        if docs:
            vectorstore = process_documents(docs)
            qa_chain = create_chain(vectorstore)
        else:
            return [[None, "⚠ Error: Could not load documents."]]

    if qa_chain is None:
        return [[None, "⚠ Please upload documents first."]]

    # Translate Urdu input to English (optional)
    if language == "Urdu":
        question = translate_from_urdu(question)

    result = qa_chain.invoke({"question": question})
    answer = result.get("answer", "I don't have that information in the provided documents.")

    # Translate back to Urdu
    if language == "Urdu":
        answer = translate_to_urdu(answer)

    return [[question, answer]]

# ------------------------
# Run UI
# ------------------------
if __name__ == "__main__":
    demo = build_interface(chatbot_response)
    demo.launch()


# 1 used for urdu translattion model and also fisrt download this then 
# Add a page number